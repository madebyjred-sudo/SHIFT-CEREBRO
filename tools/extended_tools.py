"""
Extended Tools for Shift AI Gateway
Based on various libraries:
- Unstructured: https://github.com/Unstructured-IO/unstructured
- ReportLab: https://hg.reportlab.com/hg-public/reportlab/
- python-pptx: https://github.com/scanny/python-pptx
- Pandas: https://github.com/pandas-dev/pandas
- Matplotlib: https://github.com/matplotlib/matplotlib
- OpenAI: https://github.com/openai/openai-python
- TextBlob: https://github.com/sloria/TextBlob
- qrcode: https://github.com/lincolnloop/python-qrcode
- wordcloud: https://github.com/amueller/word_cloud

Installation:
pip install unstructured reportlab python-pptx pandas matplotlib openai textblob qrcode wordcloud pillow
"""

import os
import io
import base64
import tempfile
from datetime import datetime
from typing import Optional, List, Dict, Any
from langchain_core.tools import tool

# Local storage configuration
DOCUMENTS_DIR = os.path.join(os.path.dirname(__file__), "..", "generated_documents")
os.makedirs(DOCUMENTS_DIR, exist_ok=True)


# ═══════════════════════════════════════════════════════════════
# 1. UNSTRUCTURED - Structured Document Processing
# ═══════════════════════════════════════════════════════════════

try:
    from unstructured.partition.auto import partition
    from unstructured.documents.elements import Element
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False


@tool
def generate_structured_document(
    content: str,
    document_type: str = "report",
    title: str = "Documento Estructurado",
    metadata: Optional[Dict[str, str]] = None
) -> str:
    """
    Create a structured document using Unstructured.io processing.
    
    Based on: https://github.com/Unstructured-IO/unstructured
    
    Args:
        content: Raw text content to structure
        document_type: Type of document (report, memo, contract, proposal)
        title: Document title
        metadata: Optional metadata (author, department, classification)
    
    Returns:
        str: JSON structure of the document elements
    
    Example:
        generate_structured_document(
            content="This is a report...",
            document_type="report",
            title="Q1 Analysis",
            metadata={"author": "John", "department": "Finance"}
        )
    """
    if not UNSTRUCTURED_AVAILABLE:
        return "Error: unstructured library not installed. Run: pip install unstructured"
    
    try:
        # Create temporary file for processing
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
            f.write(content)
            temp_path = f.name
        
        # Process with unstructured
        elements = partition(filename=temp_path)
        
        # Structure elements
        structured_data = {
            "title": title,
            "document_type": document_type,
            "metadata": metadata or {},
            "elements": [],
            "timestamp": datetime.now().isoformat()
        }
        
        for element in elements:
            element_dict = {
                "type": type(element).__name__,
                "text": str(element),
                "metadata": element.metadata.to_dict() if hasattr(element, 'metadata') else {}
            }
            structured_data["elements"].append(element_dict)
        
        # Clean up
        os.unlink(temp_path)
        
        return f"Structured document generated: {len(structured_data['elements'])} elements extracted"
        
    except Exception as e:
        return f"Error processing document: {str(e)}"


# ═══════════════════════════════════════════════════════════════
# 2. REPORTLAB - PDF Generation
# ═══════════════════════════════════════════════════════════════

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch


@tool
def generate_pdf_report(
    title: str,
    content: str,
    report_type: str = "general",
    sections: Optional[List[Dict[str, str]]] = None,
    include_toc: bool = False,
    filename: Optional[str] = None
) -> str:
    """
    Create a professional PDF report using ReportLab.
    
    Based on: https://hg.reportlab.com/hg-public/reportlab/
    
    Args:
        title: Report title
        content: Main report content
        report_type: Type of report (general, financial, technical, marketing)
        sections: Optional list of sections with 'heading' and 'content'
        include_toc: Whether to include table of contents
        filename: Optional custom filename (without extension)
    
    Returns:
        str: URL/path to the generated PDF file
    
    Example:
        generate_pdf_report(
            title="Financial Analysis Q1 2026",
            content="Executive summary...",
            report_type="financial",
            sections=[{"heading": "Revenue", "content": "$1M..."}]
        )
    """
    try:
        # Generate filename
        if not filename:
            filename = f"report_{report_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        final_filename = f"{filename}.pdf"
        file_path = os.path.join(DOCUMENTS_DIR, final_filename)
        
        # Create PDF document
        doc = SimpleDocTemplate(
            file_path,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        # Container for the 'Flowable' objects
        elements = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1f4e79'),
            spaceAfter=30,
            alignment=1  # Center
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#2e74b5'),
            spaceAfter=12
        )
        
        # Add title
        elements.append(Paragraph(title, title_style))
        elements.append(Spacer(1, 0.2*inch))
        
        # Add metadata
        current_date = datetime.now().strftime("%d de %B de %Y")
        elements.append(Paragraph(f"<i>Generado: {current_date}</i>", styles['Normal']))
        elements.append(Spacer(1, 0.3*inch))
        
        # Add main content
        if content:
            elements.append(Paragraph(content, styles['BodyText']))
            elements.append(Spacer(1, 0.2*inch))
        
        # Add sections
        if sections:
            for section in sections:
                heading = section.get('heading', '')
                content = section.get('content', '')
                
                if heading:
                    elements.append(Paragraph(heading, heading_style))
                    elements.append(Spacer(1, 0.1*inch))
                
                if content:
                    elements.append(Paragraph(content, styles['BodyText']))
                    elements.append(Spacer(1, 0.2*inch))
        
        # Add footer
        elements.append(Spacer(1, 0.3*inch))
        footer_style = ParagraphStyle(
            'Footer',
            parent=styles['Normal'],
            fontSize=8,
            textColor=colors.grey,
            alignment=1
        )
        elements.append(Paragraph("─  Generado por Shift AI  ─", footer_style))
        
        # Build PDF
        doc.build(elements)
        
        return f"/documents/{final_filename}"
        
    except Exception as e:
        return f"Error creating PDF: {str(e)}"


# ═══════════════════════════════════════════════════════════════
# 3. PYTHON-PPTX - PowerPoint Presentations
# ═══════════════════════════════════════════════════════════════

from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor


@tool
def create_presentation(
    title: str,
    subtitle: Optional[str] = None,
    slides_content: Optional[List[Dict[str, Any]]] = None,
    template: str = "default",
    filename: Optional[str] = None
) -> str:
    """
    Create a PowerPoint presentation using python-pptx.
    
    Based on: https://github.com/scanny/python-pptx
    
    Args:
        title: Presentation title
        subtitle: Optional subtitle
        slides_content: List of slide dictionaries with 'title' and 'content'
        template: Template style (default, corporate, minimal)
        filename: Optional custom filename (without extension)
    
    Returns:
        str: URL/path to the generated PPTX file
    
    Example:
        create_presentation(
            title="Marketing Strategy 2026",
            subtitle="Q1-Q4 Roadmap",
            slides_content=[
                {"title": "Overview", "content": "Our strategy focuses on..."},
                {"title": "Goals", "content": "1. Increase revenue by 25%"}
            ]
        )
    """
    try:
        # Create presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        if subtitle:
            subtitle_shape.text = subtitle
        
        # Add content slides
        if slides_content:
            for slide_data in slides_content:
                bullet_slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(bullet_slide_layout)
                
                shapes = slide.shapes
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = slide_data.get('title', '')
                
                content = slide_data.get('content', '')
                tf = body_shape.text_frame
                tf.text = content
                
                # Format text
                for paragraph in tf.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT
        
        # Generate filename
        if not filename:
            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        final_filename = f"{filename}.pptx"
        file_path = os.path.join(DOCUMENTS_DIR, final_filename)
        
        # Save presentation
        prs.save(file_path)
        
        return f"/documents/{final_filename}"
        
    except Exception as e:
        return f"Error creating presentation: {str(e)}"


# ═══════════════════════════════════════════════════════════════
# 4. PANDAS - Data Analysis
# ═══════════════════════════════════════════════════════════════

import pandas as pd
import json


@tool
def analyze_data_table(
    data: str,
    data_format: str = "csv",
    analysis_type: str = "summary",
    columns: Optional[List[str]] = None
) -> str:
    """
    Analyze tabular data using Pandas.
    
    Based on: https://github.com/pandas-dev/pandas
    
    Args:
        data: Data content (CSV string or JSON)
        data_format: Format of data (csv, json, tsv)
        analysis_type: Type of analysis (summary, stats, correlations, trends)
        columns: Optional list of columns to analyze
    
    Returns:
        str: Analysis results as formatted text
    
    Example:
        analyze_data_table(
            data="name,revenue\nA,100\nB,200",
            data_format="csv",
            analysis_type="summary"
        )
    """
    try:
        # Parse data based on format
        if data_format.lower() == "csv":
            from io import StringIO
            df = pd.read_csv(StringIO(data))
        elif data_format.lower() == "json":
            df = pd.read_json(data)
        elif data_format.lower() == "tsv":
            from io import StringIO
            df = pd.read_csv(StringIO(data), sep='\t')
        else:
            return f"Unsupported format: {data_format}"
        
        # Filter columns if specified
        if columns:
            df = df[[col for col in columns if col in df.columns]]
        
        # Perform analysis
        if analysis_type == "summary":
            result = df.describe().to_string()
        elif analysis_type == "stats":
            result = f"Shape: {df.shape}\nColumns: {list(df.columns)}\n\n{df.describe().to_string()}"
        elif analysis_type == "correlations":
            numeric_df = df.select_dtypes(include=[float, int])
            if not numeric_df.empty:
                result = numeric_df.corr().to_string()
            else:
                result = "No numeric columns for correlation analysis"
        elif analysis_type == "trends":
            # Simple trend analysis on first numeric column
            numeric_cols = df.select_dtypes(include=[float, int]).columns
            if len(numeric_cols) > 0:
                col = numeric_cols[0]
                trend = "increasing" if df[col].iloc[-1] > df[col].iloc[0] else "decreasing"
                result = f"Trend in '{col}': {trend}\nStart: {df[col].iloc[0]}\nEnd: {df[col].iloc[-1]}"
            else:
                result = "No numeric columns for trend analysis"
        else:
            result = df.head(10).to_string()
        
        return f"Analysis Results:\n{result}"
        
    except Exception as e:
        return f"Error analyzing data: {str(e)}"


# ═══════════════════════════════════════════════════════════════
# 5. MATPLOTLIB - Chart Visualization
# ═══════════════════════════════════════════════════════════════

import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import numpy as np


@tool
def create_chart_visualization(
    chart_type: str,
    data: str,
    title: str = "Chart",
    x_label: Optional[str] = None,
    y_label: Optional[str] = None,
    filename: Optional[str] = None
) -> str:
    """
    Create a chart visualization using Matplotlib.
    
    Based on: https://github.com/matplotlib/matplotlib
    
    Args:
        chart_type: Type of chart (bar, line, pie, scatter, histogram)
        data: JSON string with 'labels' and 'values' arrays
        title: Chart title
        x_label: X-axis label
        y_label: Y-axis label
        filename: Optional custom filename (without extension)
    
    Returns:
        str: URL/path to the generated chart image
    
    Example:
        create_chart_visualization(
            chart_type="bar",
            data='{"labels": ["A", "B", "C"], "values": [10, 20, 15]}',
            title="Sales by Region"
        )
    """
    try:
        # Parse data
        data_dict = json.loads(data)
        labels = data_dict.get('labels', [])
        values = data_dict.get('values', [])
        
        if not labels or not values:
            return "Error: Missing labels or values in data"
        
        # Create figure
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # Create chart based on type
        if chart_type == "bar":
            ax.bar(labels, values, color='#2e74b5')
        elif chart_type == "line":
            ax.plot(labels, values, marker='o', color='#2e74b5', linewidth=2)
        elif chart_type == "pie":
            ax.pie(values, labels=labels, autopct='%1.1f%%', startangle=90)
        elif chart_type == "scatter":
            ax.scatter(range(len(values)), values, color='#2e74b5', s=100)
        elif chart_type == "histogram":
            ax.hist(values, bins=10, color='#2e74b5', edgecolor='black')
        else:
            return f"Unsupported chart type: {chart_type}"
        
        # Add labels
        ax.set_title(title, fontsize=16, fontweight='bold')
        if x_label and chart_type != "pie":
            ax.set_xlabel(x_label)
        if y_label and chart_type != "pie":
            ax.set_ylabel(y_label)
        
        # Style
        ax.grid(True, alpha=0.3)
        plt.tight_layout()
        
        # Generate filename
        if not filename:
            filename = f"chart_{chart_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        final_filename = f"{filename}.png"
        file_path = os.path.join(DOCUMENTS_DIR, final_filename)
        
        # Save chart
        plt.savefig(file_path, dpi=150, bbox_inches='tight')
        plt.close()
        
        return f"/documents/{final_filename}"
        
    except Exception as e:
        return f"Error creating chart: {str(e)}"


# ═══════════════════════════════════════════════════════════════
# 6. OPENAI - DALL-E Image Generation
# ═══════════════════════════════════════════════════════════════

try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False


@tool
def generate_marketing_image(
    prompt: str,
    size: str = "1024x1024",
    quality: str = "standard",
    style: str = "vivid",
    filename: Optional[str] = None
) -> str:
    """
    Generate a marketing image using DALL-E 3 via OpenAI API.
    
    Based on: https://github.com/openai/openai-python
    API Docs: https://platform.openai.com/docs/guides/images
    
    Args:
        prompt: Detailed description of the image to generate
        size: Image size (1024x1024, 1024x1792, 1792x1024)
        quality: Image quality (standard, hd)
        style: Image style (vivid, natural)
        filename: Optional custom filename (without extension)
    
    Returns:
        str: URL/path to the generated image
    
    Example:
        generate_marketing_image(
            prompt="A professional marketing banner for a tech company...",
            size="1792x1024",
            quality="hd"
        )
    """
    if not OPENAI_AVAILABLE:
        return "Error: openai library not installed. Run: pip install openai"
    
    try:
        # Check for API key
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return "Error: OPENAI_API_KEY environment variable not set"
        
        client = OpenAI(api_key=api_key)
        
        # Generate image
        response = client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size=size,
            quality=quality,
            style=style,
            n=1,
            response_format="url"
        )
        
        # Get image URL
        image_url = response.data[0].url
        
        # Also save locally if possible
        if not filename:
            filename = f"image_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        final_filename = f"{filename}.txt"  # Store URL reference
        file_path = os.path.join(DOCUMENTS_DIR, final_filename)
        
        with open(file_path, 'w') as f:
            f.write(f"DALL-E Generated Image URL: {image_url}\nPrompt: {prompt}\n")
        
        return f"Image generated successfully. URL: {image_url}"
        
    except Exception as e:
        return f"Error generating image: {str(e)}"


# ═══════════════════════════════════════════════════════════════
# 7. TEXTBLOB - Sentiment Analysis
# ═══════════════════════════════════════════════════════════════

from textblob import TextBlob


@tool
def analyze_content_sentiment(
    text: str,
    analysis_type: str = "sentiment"
) -> str:
    """
    Analyze text sentiment using TextBlob.
    
    Based on: https://github.com/sloria/TextBlob
    
    Args:
        text: Text content to analyze
        analysis_type: Type of analysis (sentiment, subjectivity, noun_phrases)
    
    Returns:
        str: Analysis results
    
    Example:
        analyze_content_sentiment(
            text="This product is amazing! I love it.",
            analysis_type="sentiment"
        )
    """
    try:
        blob = TextBlob(text)
        
        if analysis_type == "sentiment":
            polarity = blob.sentiment.polarity
            subjectivity = blob.sentiment.subjectivity
            
            if polarity > 0.1:
                sentiment = "Positive"
            elif polarity < -0.1:
                sentiment = "Negative"
            else:
                sentiment = "Neutral"
            
            result = (
                f"Sentiment: {sentiment}\n"
                f"Polarity: {polarity:.3f} (-1 to 1)\n"
                f"Subjectivity: {subjectivity:.3f} (0 to 1)"
            )
        
        elif analysis_type == "subjectivity":
            result = f"Subjectivity: {blob.sentiment.subjectivity:.3f} (0=objective, 1=subjective)"
        
        elif analysis_type == "noun_phrases":
            phrases = list(blob.noun_phrases)
            result = f"Noun Phrases: {', '.join(phrases)}"
        
        else:
            result = f"Sentiment: {blob.sentiment}"
        
        return result
        
    except Exception as e:
        return f"Error analyzing sentiment: {str(e)}"


# ═══════════════════════════════════════════════════════════════
# 8. QRCODE - QR Code Generation
# ═══════════════════════════════════════════════════════════════

import qrcode


@tool
def generate_campaign_qr(
    data: str,
    qr_type: str = "url",
    size: int = 10,
    filename: Optional[str] = None
) -> str:
    """
    Generate a QR code for campaigns using python-qrcode.
    
    Based on: https://github.com/lincolnloop/python-qrcode
    
    Args:
        data: Data to encode (URL, text, etc.)
        qr_type: Type of QR (url, text, vcard, email)
        size: QR code size (box size in pixels)
        filename: Optional custom filename (without extension)
    
    Returns:
        str: URL/path to the generated QR code image
    
    Example:
        generate_campaign_qr(
            data="https://shiftpn.com/campaign",
            qr_type="url",
            size=10
        )
    """
    try:
        # Create QR code
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_H,
            box_size=size,
            border=4,
        )
        
        # Add data
        if qr_type == "url":
            qr.add_data(data)
        elif qr_type == "email":
            qr.add_data(f"mailto:{data}")
        else:
            qr.add_data(data)
        
        qr.make(fit=True)
        
        # Create image
        img = qr.make_image(fill_color="black", back_color="white")
        
        # Generate filename
        if not filename:
            filename = f"qr_{qr_type}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        final_filename = f"{filename}.png"
        file_path = os.path.join(DOCUMENTS_DIR, final_filename)
        
        # Save image
        img.save(file_path)
        
        return f"/documents/{final_filename}"
        
    except Exception as e:
        return f"Error generating QR code: {str(e)}"


# ═══════════════════════════════════════════════════════════════
# 9. WORDCLOUD - Keyword Cloud Generation
# ═══════════════════════════════════════════════════════════════

from wordcloud import WordCloud


@tool
def generate_keyword_cloud(
    text: str,
    max_words: int = 100,
    width: int = 800,
    height: int = 400,
    background_color: str = "white",
    filename: Optional[str] = None
) -> str:
    """
    Generate a keyword cloud (word cloud) using wordcloud library.
    
    Based on: https://github.com/amueller/word_cloud
    
    Args:
        text: Text content to analyze
        max_words: Maximum number of words to display
        width: Image width in pixels
        height: Image height in pixels
        background_color: Background color (white, black, transparent)
        filename: Optional custom filename (without extension)
    
    Returns:
        str: URL/path to the generated word cloud image
    
    Example:
        generate_keyword_cloud(
            text="Marketing strategy digital content social media...",
            max_words=50,
            width=1200,
            height=600
        )
    """
    try:
        # Generate word cloud
        wordcloud = WordCloud(
            width=width,
            height=height,
            max_words=max_words,
            background_color=background_color,
            colormap='viridis',
            contour_width=1,
            contour_color='steelblue'
        ).generate(text)
        
        # Generate filename
        if not filename:
            filename = f"wordcloud_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        
        final_filename = f"{filename}.png"
        file_path = os.path.join(DOCUMENTS_DIR, final_filename)
        
        # Save image
        wordcloud.to_file(file_path)
        
        return f"/documents/{final_filename}"
        
    except Exception as e:
        return f"Error generating word cloud: {str(e)}"


# ═══════════════════════════════════════════════════════════════
# EXPORT ALL EXTENDED TOOLS
# ═══════════════════════════════════════════════════════════════

EXTENDED_TOOLS = [
    generate_structured_document,
    generate_pdf_report,
    create_presentation,
    analyze_data_table,
    create_chart_visualization,
    generate_marketing_image,
    analyze_content_sentiment,
    generate_campaign_qr,
    generate_keyword_cloud,
]